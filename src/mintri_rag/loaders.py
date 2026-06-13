from pathlib import Path

from pypdf import PdfReader


SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".txt", ".md"}


def discover_documents(raw_dir: Path) -> list[Path]:
    return sorted(
        path
        for path in raw_dir.glob("**/*")
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS
    )


def load_document(path: Path) -> str:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return _load_pdf(path)
    if suffix == ".pptx":
        return _load_pptx(path)
    if suffix in {".txt", ".md"}:
        return path.read_text(encoding="utf-8", errors="ignore")
    raise ValueError(f"Unsupported document type: {path.suffix}")


def _load_pdf(path: Path) -> str:
    reader = PdfReader(str(path))
    pages = []
    for page_number, page in enumerate(reader.pages, start=1):
        text = page.extract_text() or ""
        if text.strip():
            pages.append(f"\n[Page {page_number}]\n{text}")
    return "\n".join(pages)


def _load_pptx(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(str(path))
    slides = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        if texts:
            slides.append(f"\n[Slide {slide_number}]\n" + "\n".join(texts))
    return "\n".join(slides)
